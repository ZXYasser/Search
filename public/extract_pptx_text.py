from pptx import Presentation
from pathlib import Path
import sys

def extract(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    parts.append(t)
        if parts:
            lines.append(f"[شريحة {i}]\n" + "\n".join(parts))
    return "\n\n".join(lines)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx_text.py input.pptx output.txt")
        sys.exit(1)

    inp = sys.argv[1]
    out = sys.argv[2]
    text = extract(inp)
    Path(out).write_text(text, encoding="utf-8")
    print("Saved:", out)
